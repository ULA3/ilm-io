"""
Agentic pipeline router — educator mode.

POST /agents/read              Reader Agent  → pockets JSON
GET  /agents/pockets/{file_id} Fetch stored pockets
POST /agents/adhd-slides       Slide Sorter (ADHD)   → PPTX + slide JSON
POST /agents/autism-slides     Slide Sorter (Autism) → PPTX + slide JSON
POST /agents/transcribe        Transcriber           → audio script JSON
"""
from __future__ import annotations
import io
import math as _math
import uuid

from fastapi import APIRouter, HTTPException
from fastapi.concurrency import run_in_threadpool
from pydantic import BaseModel

from services import text_cache, artifact_cache, pocket_cache, agent_service, tts_service

router = APIRouter(prefix="/agents", tags=["agents"])


# ── Request models ─────────────────────────────────────────────────────────

class ReadRequest(BaseModel):
    file_id: str
    lang: str = "en"     # en | ms | zh | ta


class AgentGenerateRequest(BaseModel):
    file_id: str
    lang: str = "en"     # en | ms | zh | ta


# ── PPTX builders ──────────────────────────────────────────────────────────

_THEME_COLORS = {
    "teal":   {"bg": (0xE0, 0xF2, 0xF1), "accent": (0x00, 0x96, 0x88), "text": (0x00, 0x4D, 0x40)},
    "amber":  {"bg": (0xFF, 0xF8, 0xE1), "accent": (0xFF, 0xB3, 0x00), "text": (0x4E, 0x34, 0x00)},
    "rose":   {"bg": (0xFF, 0xEB, 0xEE), "accent": (0xE5, 0x39, 0x35), "text": (0x4A, 0x00, 0x02)},
    "violet": {"bg": (0xED, 0xE7, 0xF6), "accent": (0x73, 0x44, 0xB8), "text": (0x31, 0x00, 0x77)},
}
_DEFAULT_THEME = _THEME_COLORS["teal"]


def _build_adhd_pptx(slides: list[dict], topic: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for s in slides:
        theme_name = s.get("color_theme", "teal")
        th = _THEME_COLORS.get(theme_name, _DEFAULT_THEME)
        BG     = RGBColor(*th["bg"])
        ACCENT = RGBColor(*th["accent"])
        TEXT   = RGBColor(*th["text"])
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG

        # ── Accent bar (left edge) ──
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        # ── Timer badge (top-right) ──
        timer = s.get("timer_minutes")
        if timer:
            tb = slide.shapes.add_textbox(Inches(11.3), Inches(0.25), Inches(1.7), Inches(0.55))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = f"⏱  {timer} min"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.alignment = PP_ALIGN.RIGHT

        # ── Slide index (top-left after bar) ──
        idx_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(0.8), Inches(0.5))
        ip = idx_box.text_frame.paragraphs[0]
        ip.text = str(s.get("index", ""))
        ip.font.size = Pt(13)
        ip.font.color.rgb = ACCENT

        # ── Title ──
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(10.5), Inches(1.3))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s.get("title", "")
        p.font.bold = True
        p.font.size = Pt(38)
        p.font.color.rgb = TEXT

        # ── Bullets ──
        bullet_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(7.8), Inches(3.5))
        tf2 = bullet_box.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(s.get("bullets", [])[:3]):
            para = tf2.add_paragraph() if i else tf2.paragraphs[0]
            para.text = f"→  {bullet}"
            para.font.size = Pt(24)
            para.font.color.rgb = TEXT
            para.space_before = Pt(10)

        # ── Visual hint (right column) ──
        hint = s.get("visual_hint", "")
        if hint:
            vh = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(4.0), Inches(3.5))
            vf = vh.text_frame
            vf.word_wrap = True
            vp = vf.paragraphs[0]
            vp.text = "🖼  " + hint
            vp.font.size = Pt(13)
            vp.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
            vp.font.italic = True

        # ── Focus question box (bottom) ──
        fq = s.get("focus_question", "")
        if fq:
            fqb = slide.shapes.add_textbox(Inches(0.6), Inches(5.9), Inches(12.0), Inches(0.8))
            fqb_fill = fqb.fill
            fqb_fill.solid()
            fqb_fill.fore_color.rgb = ACCENT
            fqf = fqb.text_frame
            fqp = fqf.paragraphs[0]
            fqp.text = f"🎯  {fq}"
            fqp.font.size = Pt(16)
            fqp.font.bold = True
            fqp.font.color.rgb = WHITE

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_autism_pptx(slides: list[dict], topic: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Fixed palette — high contrast, zero decoration
    BG      = RGBColor(0xFF, 0xFF, 0xFF)   # white
    HEADER  = RGBColor(0x1A, 0x23, 0x7E)   # deep navy
    LABEL   = RGBColor(0x1A, 0x23, 0x7E)   # same navy for labels
    BODY    = RGBColor(0x1A, 0x1A, 0x1A)   # near-black
    BORDER  = RGBColor(0xE0, 0xE0, 0xE0)   # light gray divider
    NUM_BG  = RGBColor(0x1A, 0x23, 0x7E)   # navy for slide number chip

    for s in slides:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG

        # ── Top navy bar ──
        top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.4))
        top.fill.solid()
        top.fill.fore_color.rgb = HEADER
        top.line.fill.background()

        # ── Slide number chip ──
        ni = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(0.7), Inches(0.7))
        nf = ni.text_frame
        np = nf.paragraphs[0]
        np.text = str(s.get("index", ""))
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # ── Heading ──
        hb = slide.shapes.add_textbox(Inches(1.1), Inches(0.2), Inches(11.8), Inches(1.0))
        hf = hb.text_frame
        hp = hf.paragraphs[0]
        hp.text = s.get("heading", "")
        hp.font.size = Pt(28)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        hf.word_wrap = True

        # ── "What:" section ──
        wl = slide.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(1.2), Inches(0.4))
        wl.text_frame.paragraphs[0].text = "What:"
        wl.text_frame.paragraphs[0].font.bold = True
        wl.text_frame.paragraphs[0].font.size = Pt(14)
        wl.text_frame.paragraphs[0].font.color.rgb = LABEL

        wb = slide.shapes.add_textbox(Inches(1.75), Inches(1.6), Inches(8.5), Inches(0.6))
        wf = wb.text_frame
        wf.word_wrap = True
        wp = wf.paragraphs[0]
        wp.text = s.get("what", "")
        wp.font.size = Pt(16)
        wp.font.color.rgb = BODY

        # ── Divider ──
        ln = slide.shapes.add_shape(1, Inches(0.5), Inches(2.35), Inches(12.3), Inches(0.02))
        ln.fill.solid()
        ln.fill.fore_color.rgb = BORDER
        ln.line.fill.background()

        # ── "Details:" section ──
        dl = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(1.4), Inches(0.4))
        dl.text_frame.paragraphs[0].text = "Details:"
        dl.text_frame.paragraphs[0].font.bold = True
        dl.text_frame.paragraphs[0].font.size = Pt(14)
        dl.text_frame.paragraphs[0].font.color.rgb = LABEL

        detail_box = slide.shapes.add_textbox(Inches(1.95), Inches(2.45), Inches(7.5), Inches(2.6))
        df = detail_box.text_frame
        df.word_wrap = True
        for i, detail in enumerate(s.get("details", [])[:4]):
            para = df.add_paragraph() if i else df.paragraphs[0]
            para.text = f"{i + 1}.  {detail}"
            para.font.size = Pt(16)
            para.font.color.rgb = BODY
            para.space_before = Pt(6)

        # ── "Why it matters:" section ──
        wml = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(2.3), Inches(0.4))
        wml.text_frame.paragraphs[0].text = "Why it matters:"
        wml.text_frame.paragraphs[0].font.bold = True
        wml.text_frame.paragraphs[0].font.size = Pt(13)
        wml.text_frame.paragraphs[0].font.color.rgb = LABEL

        wmb = slide.shapes.add_textbox(Inches(2.9), Inches(5.15), Inches(10.0), Inches(0.6))
        wmf = wmb.text_frame
        wmf.word_wrap = True
        wmp = wmf.paragraphs[0]
        wmp.text = s.get("why_it_matters", "")
        wmp.font.size = Pt(14)
        wmp.font.italic = True
        wmp.font.color.rgb = BODY

        # ── Visual description (bottom-right) ──
        vd = s.get("visual_description", "")
        if vd:
            vb = slide.shapes.add_textbox(Inches(9.5), Inches(2.5), Inches(3.5), Inches(2.5))
            vf = vb.text_frame
            vf.word_wrap = True
            vp = vf.paragraphs[0]
            vp.text = "[ Diagram ]\n" + vd
            vp.font.size = Pt(11)
            vp.font.color.rgb = RGBColor(0x75, 0x75, 0x75)
            vp.font.italic = True

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Helper: build a flat script string from Transcriber output ─────────────

def _flatten_script(script: dict) -> str:
    parts = [script.get("intro", "")]
    for seg in script.get("segments", []):
        parts.append(seg.get("narrative", ""))
        t = seg.get("transition", "")
        if t:
            parts.append(t)
    parts.append(script.get("outro", ""))
    return "\n\n".join(p for p in parts if p)


# ── Endpoints ──────────────────────────────────────────────────────────────

@router.post("/read")
async def read_document(req: ReadRequest):
    """
    Reader Agent — step 1 of the pipeline.
    Extracts structured learning pockets from the uploaded file.
    Result is cached by file_id so subsequent generate calls are fast.
    """
    cached = pocket_cache.get(req.file_id)
    if cached:
        return cached

    text = text_cache.get(req.file_id)
    if not text:
        raise HTTPException(404, f"File '{req.file_id}' not found. Please re-upload.")

    pockets = await run_in_threadpool(agent_service.read_document, text, req.lang)
    pocket_cache.store(req.file_id, pockets)
    return pockets


@router.get("/pockets/{file_id}")
async def get_pockets(file_id: str):
    """Return previously extracted pockets (from cache)."""
    pockets = pocket_cache.get(file_id)
    if not pockets:
        raise HTTPException(404, "Pockets not found. Call POST /agents/read first.")
    return pockets


@router.post("/adhd-slides")
async def generate_adhd_slides(req: AgentGenerateRequest):
    """
    Slide Sorter (ADHD) — generates an ADHD-optimised PPTX from pockets.
    Uses the Orchestrator's pocket cache; re-reads the document if pockets are absent.
    """
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    job_id = str(uuid.uuid4())
    slides = await run_in_threadpool(agent_service.generate_adhd_slides, pockets, req.lang)
    topic  = pockets.get("topic", "Slides")
    pptx   = await run_in_threadpool(_build_adhd_pptx, slides, topic)

    artifact_cache.store(
        job_id, pptx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        f"adhd_slides_{job_id[:8]}.pptx",
    )

    return {
        "job_id":       job_id,
        "slides":       slides,
        "download_url": f"/api/download/{job_id}",
        "topic":        topic,
        "slide_count":  len(slides),
    }


@router.post("/autism-slides")
async def generate_autism_slides(req: AgentGenerateRequest):
    """
    Slide Sorter (Autism) — generates a structured, predictable PPTX from pockets.
    """
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    job_id = str(uuid.uuid4())
    slides = await run_in_threadpool(agent_service.generate_autism_slides, pockets, req.lang)
    topic  = pockets.get("topic", "Slides")
    pptx   = await run_in_threadpool(_build_autism_pptx, slides, topic)

    artifact_cache.store(
        job_id, pptx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        f"autism_slides_{job_id[:8]}.pptx",
    )

    return {
        "job_id":       job_id,
        "slides":       slides,
        "download_url": f"/api/download/{job_id}",
        "topic":        topic,
        "slide_count":  len(slides),
    }


@router.post("/transcribe")
async def transcribe(req: AgentGenerateRequest):
    """
    Transcriber Agent — generates a flowing audio narrative script from pockets
    and renders it to an MP3 via TTS. Returns script JSON + audio download URL.
    """
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    script    = await run_in_threadpool(agent_service.transcribe_to_audio, pockets, req.lang)
    full_text = _flatten_script(script)

    # Generate MP3 and cache it
    audio_url = None
    try:
        mp3_bytes = await tts_service.text_to_speech(full_text)
        audio_job_id = str(uuid.uuid4())
        artifact_cache.store(audio_job_id, mp3_bytes, "audio/mpeg", f"audiobook_{audio_job_id[:8]}.mp3")
        audio_url = f"/api/download/{audio_job_id}"
    except Exception:
        pass  # audio optional — still return the script

    return {
        "script":                     script,
        "full_text":                  full_text,
        "audio_url":                  audio_url,
        "estimated_duration_minutes": script.get("estimated_duration_minutes", 5),
        "word_count":                 script.get("word_count", len(full_text.split())),
        "topic":                      pockets.get("topic", ""),
    }


# ── Additional request models ──────────────────────────────────────────────

class CuriousCriticRequest(BaseModel):
    file_id: str
    condition: str = "general"   # adhd | dyslexia | autism | general
    lang: str = "en"


class IlmuistRequest(BaseModel):
    message: str
    history: list[dict] = []
    file_id: str | None = None
    role: str = "teacher"   # "teacher" | "student"


class StudentChatActionRequest(BaseModel):
    file_id: str
    action: str   # summarize | quiz | vocab | study_plan | focus_tip
    lang: str = "en"
    mood: str = "okay"


class StudentObserverRequest(BaseModel):
    students: list[dict]
    session_events: list[dict] = []
    topic: str = ""


# ── Dyslexia PPTX builder ──────────────────────────────────────────────────

def _build_dyslexia_pptx(slides: list[dict], topic: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    import lxml.etree as etree

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    OFFWHITE  = RGBColor(0xFA, 0xF8, 0xF5)   # warm off-white — easier on eyes
    DARK      = RGBColor(0x1C, 0x1C, 0x1C)   # near-black text
    ACCENT    = RGBColor(0x1A, 0x5C, 0x96)   # accessible blue for titles
    YELLOW_BG = RGBColor(0xFF, 0xF0, 0x80)   # soft yellow for key phrase
    NOTE_BG   = RGBColor(0xE8, 0xF4, 0xFD)   # pale blue reading note strip

    def _set_line_spacing(text_frame, spacing_pt: float):
        for para in text_frame.paragraphs:
            pPr = para._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            spcPts = etree.SubElement(lnSpc, qn("a:spcPts"))
            spcPts.set("val", str(int(spacing_pt * 100)))

    for s in slides:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = OFFWHITE

        # ── Accent top strip ──
        top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.18))
        top.fill.solid()
        top.fill.fore_color.rgb = ACCENT
        top.line.fill.background()

        # ── Slide number (top-left) ──
        ni = slide.shapes.add_textbox(Inches(0.3), Inches(0.28), Inches(0.5), Inches(0.4))
        np_ = ni.text_frame.paragraphs[0]
        np_.text = str(s.get("index", ""))
        np_.font.size = Pt(13)
        np_.font.color.rgb = ACCENT

        # ── Title — large, left-aligned, bold ──
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tp = tf.paragraphs[0]
        tp.text = s.get("title", "")
        tp.font.bold = True
        tp.font.size = Pt(32)
        tp.font.color.rgb = ACCENT
        tp.font.name = "Arial"
        tp.alignment = PP_ALIGN.LEFT

        # ── Body bullets — 18pt Arial, left-aligned, generous spacing ──
        bb = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(8.8), Inches(3.8))
        bf = bb.text_frame
        bf.word_wrap = True
        for i, bullet in enumerate(s.get("bullets", [])[:4]):
            para = bf.add_paragraph() if i else bf.paragraphs[0]
            para.text = bullet
            para.font.size = Pt(18)
            para.font.color.rgb = DARK
            para.font.name = "Arial"
            para.alignment = PP_ALIGN.LEFT
            para.space_before = Pt(8)
            para.space_after = Pt(6)
        _set_line_spacing(bf, 27)   # ~1.8× line height for 18pt text

        # ── Key phrase highlight box (right column) ──
        kp = s.get("key_phrase", "")
        if kp:
            kpb = slide.shapes.add_textbox(Inches(9.6), Inches(1.7), Inches(3.4), Inches(1.0))
            kpb.fill.solid()
            kpb.fill.fore_color.rgb = YELLOW_BG
            kpf = kpb.text_frame
            kpf.word_wrap = True
            kp_label = kpf.paragraphs[0]
            kp_label.text = "Key phrase:"
            kp_label.font.size = Pt(11)
            kp_label.font.bold = True
            kp_label.font.color.rgb = DARK
            kp_label.font.name = "Arial"
            kp2 = kpf.add_paragraph()
            kp2.text = kp
            kp2.font.size = Pt(16)
            kp2.font.bold = True
            kp2.font.color.rgb = DARK
            kp2.font.name = "Arial"

        # ── Visual hint ──
        vh = s.get("visual_hint", "")
        if vh:
            vhb = slide.shapes.add_textbox(Inches(9.6), Inches(3.0), Inches(3.4), Inches(2.2))
            vhf = vhb.text_frame
            vhf.word_wrap = True
            vhp = vhf.paragraphs[0]
            vhp.text = "[ Image ]\n" + vh
            vhp.font.size = Pt(11)
            vhp.font.italic = True
            vhp.font.color.rgb = RGBColor(0x70, 0x70, 0x70)
            vhp.font.name = "Arial"

        # ── Reading note strip (bottom) ──
        rn = s.get("reading_note", "")
        if rn:
            rnb = slide.shapes.add_textbox(Inches(0), Inches(6.5), Inches(13.33), Inches(0.7))
            rnb.fill.solid()
            rnb.fill.fore_color.rgb = NOTE_BG
            rnf = rnb.text_frame
            rnp = rnf.paragraphs[0]
            rnp.text = "📖  " + rn
            rnp.font.size = Pt(13)
            rnp.font.color.rgb = RGBColor(0x1A, 0x5C, 0x96)
            rnp.font.name = "Arial"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── New endpoints ──────────────────────────────────────────────────────────

@router.post("/dyslexia-slides")
async def generate_dyslexia_slides(req: AgentGenerateRequest):
    """
    Slide Sorter (Dyslexia) — generates a dyslexia-friendly PPTX from pockets.
    Arial font, off-white background, generous spacing, key phrase highlights.
    """
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    job_id = str(uuid.uuid4())
    slides = await run_in_threadpool(agent_service.generate_dyslexia_slides, pockets, req.lang)
    topic  = pockets.get("topic", "Slides")
    pptx   = await run_in_threadpool(_build_dyslexia_pptx, slides, topic)

    artifact_cache.store(
        job_id, pptx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        f"dyslexia_slides_{job_id[:8]}.pptx",
    )

    return {
        "job_id":       job_id,
        "slides":       slides,
        "download_url": f"/api/download/{job_id}",
        "topic":        topic,
        "slide_count":  len(slides),
    }


@router.post("/worksheet")
async def generate_worksheet(req: CuriousCriticRequest):
    """
    Curious Critic — generates a condition-tailored worksheet from pockets.
    """
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    worksheet = await run_in_threadpool(
        agent_service.curious_critic_worksheet, pockets, req.condition, req.lang
    )
    return {
        "worksheet":  worksheet,
        "condition":  req.condition,
        "topic":      pockets.get("topic", ""),
    }


@router.post("/quiz")
async def generate_quiz(req: CuriousCriticRequest):
    """
    Curious Critic — generates a condition-tailored quiz from pockets.
    """
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    quiz = await run_in_threadpool(
        agent_service.curious_critic_quiz, pockets, req.condition, req.lang
    )
    return {
        "quiz":       quiz,
        "condition":  req.condition,
        "topic":      pockets.get("topic", ""),
    }


@router.post("/student-observer")
async def student_observer(req: StudentObserverRequest):
    """
    Student Observer — analyses per-student learning patterns and flags
    students who need attention. Expects a list of student profile dicts.
    """
    analysis = await run_in_threadpool(
        agent_service.student_observer_analyze,
        req.students,
        req.session_events,
        req.topic,
    )
    return {"students": analysis, "topic": req.topic}


@router.post("/student-chat-action")
async def student_chat_action(req: StudentChatActionRequest):
    """Functional Ilm actions — summarize, quiz, vocab, study plan, focus tip."""
    from services import text_cache, pocket_cache

    text = text_cache.get(req.file_id)
    pockets = pocket_cache.get(req.file_id)
    if not text:
        from services.supabase_client import get_upload_sync
        record = await run_in_threadpool(get_upload_sync, req.file_id)
        if not record:
            raise HTTPException(404, "Upload not found — please upload your file again.")
        text = record.get("text_content", "")

    result = await run_in_threadpool(
        agent_service.student_chat_action,
        req.action,
        text,
        pockets,
        req.lang,
        req.mood,
    )
    return result


@router.post("/ilmuist")
async def ilmuist_chat(req: IlmuistRequest):
    """
    Ilmuist — context-aware Manglish AI guide.
    Optionally enriches context with pocket_cache data if file_id provided.
    """
    context: dict = {}

    if req.file_id:
        pockets = pocket_cache.get(req.file_id)
        if pockets:
            pocket_list = pockets.get("pockets", [])
            context["topic"]        = pockets.get("topic", "")
            context["summary"]      = pockets.get("summary", "")
            context["pocket_count"] = len(pocket_list)

    response = await run_in_threadpool(
        agent_service.ilmuist_chat, req.message, req.history, context, req.role
    )
    return response


# ═══════════════════════════════════════════════════════════════════════════
# STUDENT PIPELINE — bionic reading helpers + PPTX builders + endpoints
# ═══════════════════════════════════════════════════════════════════════════

# ── Bionic reading helper ──────────────────────────────────────────────────

def _bionic_runs(paragraph, text: str, font_size_pt: float, color_rgb, space_before_pt: float = 0):
    """Fill `paragraph` with bionic-reading runs (first ~45% of each word bolded)."""
    from pptx.util import Pt
    if space_before_pt:
        paragraph.space_before = Pt(space_before_pt)
    words = text.split(" ")
    for i, word in enumerate(words):
        if not word:
            continue
        n = max(1, _math.ceil(len(word) * 0.45))
        bold_part, rest = word[:n], word[n:]
        for txt, bold in [(bold_part, True), (rest, False)]:
            if not txt:
                continue
            run = paragraph.add_run()
            run.text = txt
            run.font.bold = bold
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = color_rgb
        if i < len(words) - 1:
            sp = paragraph.add_run()
            sp.text = " "
            sp.font.bold = False
            sp.font.size = Pt(font_size_pt)
            sp.font.color.rgb = color_rgb


# ── Student ADHD PPTX builder ──────────────────────────────────────────────

def _build_student_adhd_pptx(slides: list[dict], topic: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for s in slides:
        theme_name = s.get("color_theme", "teal")
        th = _THEME_COLORS.get(theme_name, _DEFAULT_THEME)
        BG     = RGBColor(*th["bg"])
        ACCENT = RGBColor(*th["accent"])
        TEXT   = RGBColor(*th["text"])
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
        HONEY  = RGBColor(0xFF, 0xB3, 0x00)

        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG

        # Accent left bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

        # Timer badge
        timer = s.get("timer_minutes")
        if timer:
            tb = slide.shapes.add_textbox(Inches(11.4), Inches(0.2), Inches(1.6), Inches(0.5))
            p = tb.text_frame.paragraphs[0]
            p.text = f"⏱  {timer} min"; p.font.size = Pt(13); p.font.bold = True
            p.font.color.rgb = ACCENT; p.alignment = PP_ALIGN.RIGHT

        # Slide index chip
        idx_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.2), Inches(0.6), Inches(0.4))
        ip = idx_box.text_frame.paragraphs[0]
        ip.text = str(s.get("index", "")); ip.font.size = Pt(12); ip.font.color.rgb = ACCENT

        # Title — bionic reading
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(10.5), Inches(1.15))
        tf = tb2.text_frame; tf.word_wrap = True
        para = tf.paragraphs[0]
        _bionic_runs(para, s.get("title", ""), 34, TEXT)

        # Bullets — bionic reading
        bb = slide.shapes.add_textbox(Inches(0.65), Inches(1.85), Inches(7.8), Inches(3.2))
        bf = bb.text_frame; bf.word_wrap = True
        for i, bullet in enumerate(s.get("bullets", [])[:3]):
            para = bf.add_paragraph() if i else bf.paragraphs[0]
            _bionic_runs(para, f"→  {bullet}", 22, TEXT, space_before_pt=9 if i else 0)

        # Mind map hint (right column)
        mmh = s.get("mind_map_hint", "")
        if mmh:
            mhb = slide.shapes.add_textbox(Inches(9.1), Inches(1.85), Inches(3.9), Inches(1.0))
            mhb.fill.solid(); mhb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            mhf = mhb.text_frame; mhf.word_wrap = True
            lbl = mhf.paragraphs[0]
            lbl.text = "🗺  Map"; lbl.font.size = Pt(10); lbl.font.bold = True; lbl.font.color.rgb = ACCENT
            _bionic_runs(mhf.add_paragraph(), mmh, 12, TEXT)

        # Visual hint (right column below map)
        vh = s.get("visual_hint", "")
        if vh:
            vhb = slide.shapes.add_textbox(Inches(9.1), Inches(3.05), Inches(3.9), Inches(1.8))
            vhf = vhb.text_frame; vhf.word_wrap = True
            vhp = vhf.paragraphs[0]
            vhp.text = "🖼  " + vh; vhp.font.size = Pt(12); vhp.font.italic = True
            vhp.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

        # Focus question box
        fq = s.get("focus_question", "")
        if fq:
            fqb = slide.shapes.add_textbox(Inches(0.5), Inches(5.15), Inches(8.2), Inches(0.7))
            fqb.fill.solid(); fqb.fill.fore_color.rgb = ACCENT
            fqf = fqb.text_frame
            fqp = fqf.paragraphs[0]
            _bionic_runs(fqp, f"🎯  {fq}", 15, WHITE)

        # Fun fact strip at bottom
        ff = s.get("fun_fact", "")
        if ff:
            ffb = slide.shapes.add_textbox(Inches(0), Inches(6.3), Inches(13.33), Inches(0.9))
            ffb.fill.solid(); ffb.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xC0)
            fff = ffb.text_frame; fff.word_wrap = True
            ffp = fff.paragraphs[0]
            _bionic_runs(ffp, f"⚡  {ff}", 14, RGBColor(0x4E, 0x34, 0x00))

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


# ── Student Autism PPTX builder ────────────────────────────────────────────

def _build_student_autism_pptx(slides: list[dict], topic: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    BG     = RGBColor(0xFF, 0xFF, 0xFF)
    HEADER = RGBColor(0x1A, 0x23, 0x7E)
    LABEL  = RGBColor(0x1A, 0x23, 0x7E)
    BODY   = RGBColor(0x1A, 0x1A, 0x1A)
    BORDER = RGBColor(0xE0, 0xE0, 0xE0)

    for s in slides:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = BG

        # Top navy bar
        top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.4))
        top.fill.solid(); top.fill.fore_color.rgb = HEADER; top.line.fill.background()

        # Slide number
        ni = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(0.7), Inches(0.7))
        np_ = ni.text_frame.paragraphs[0]
        np_.text = str(s.get("index", "")); np_.font.size = Pt(18)
        np_.font.bold = True; np_.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Heading — bionic
        hb = slide.shapes.add_textbox(Inches(1.1), Inches(0.2), Inches(11.8), Inches(1.0))
        hf = hb.text_frame; hf.word_wrap = True
        _bionic_runs(hf.paragraphs[0], s.get("heading", ""), 26, RGBColor(0xFF, 0xFF, 0xFF))

        # What — bionic
        wl = slide.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(1.2), Inches(0.4))
        wl.text_frame.paragraphs[0].text = "What:"
        wl.text_frame.paragraphs[0].font.bold = True
        wl.text_frame.paragraphs[0].font.size = Pt(14)
        wl.text_frame.paragraphs[0].font.color.rgb = LABEL
        wb = slide.shapes.add_textbox(Inches(1.75), Inches(1.6), Inches(8.5), Inches(0.7))
        wb.text_frame.word_wrap = True
        _bionic_runs(wb.text_frame.paragraphs[0], s.get("what", ""), 16, BODY)

        # Divider
        ln = slide.shapes.add_shape(1, Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.02))
        ln.fill.solid(); ln.fill.fore_color.rgb = BORDER; ln.line.fill.background()

        # Details — bionic
        dl = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(1.4), Inches(0.4))
        dl.text_frame.paragraphs[0].text = "Details:"
        dl.text_frame.paragraphs[0].font.bold = True
        dl.text_frame.paragraphs[0].font.size = Pt(14)
        dl.text_frame.paragraphs[0].font.color.rgb = LABEL
        detail_box = slide.shapes.add_textbox(Inches(1.95), Inches(2.55), Inches(7.5), Inches(2.4))
        df = detail_box.text_frame; df.word_wrap = True
        for i, detail in enumerate(s.get("details", [])[:3]):
            para = df.add_paragraph() if i else df.paragraphs[0]
            _bionic_runs(para, f"{i + 1}.  {detail}", 16, BODY, space_before_pt=6 if i else 0)

        # Why it matters — bionic
        wml = slide.shapes.add_textbox(Inches(0.5), Inches(5.15), Inches(2.3), Inches(0.4))
        wml.text_frame.paragraphs[0].text = "Why it matters:"
        wml.text_frame.paragraphs[0].font.bold = True
        wml.text_frame.paragraphs[0].font.size = Pt(13)
        wml.text_frame.paragraphs[0].font.color.rgb = LABEL
        wmb = slide.shapes.add_textbox(Inches(2.9), Inches(5.1), Inches(10.0), Inches(0.7))
        wmb.text_frame.word_wrap = True
        _bionic_runs(wmb.text_frame.paragraphs[0], s.get("why_it_matters", ""), 14, BODY)

        # Visual description
        vd = s.get("visual_description", "")
        if vd:
            vb = slide.shapes.add_textbox(Inches(9.6), Inches(2.55), Inches(3.4), Inches(2.2))
            vf = vb.text_frame; vf.word_wrap = True
            vp = vf.paragraphs[0]
            vp.text = "[ Diagram ]\n" + vd
            vp.font.size = Pt(11); vp.font.italic = True
            vp.font.color.rgb = RGBColor(0x75, 0x75, 0x75)

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


# ── Mind Map PPTX builder ──────────────────────────────────────────────────

def _build_mindmap_pptx(mindmap: dict, topic: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    BRANCH_COLORS = {
        "teal":   (RGBColor(0x00, 0x96, 0x88), RGBColor(0xE0, 0xF2, 0xF1), RGBColor(0x00, 0x4D, 0x40)),
        "amber":  (RGBColor(0xFF, 0xB3, 0x00), RGBColor(0xFF, 0xF8, 0xE1), RGBColor(0x4E, 0x34, 0x00)),
        "rose":   (RGBColor(0xE5, 0x39, 0x35), RGBColor(0xFF, 0xEB, 0xEE), RGBColor(0x4A, 0x00, 0x02)),
        "violet": (RGBColor(0x73, 0x44, 0xB8), RGBColor(0xED, 0xE7, 0xF6), RGBColor(0x31, 0x00, 0x77)),
        "sage":   (RGBColor(0x5E, 0x7D, 0x5C), RGBColor(0xDC, 0xE8, 0xDB), RGBColor(0x1B, 0x34, 0x1A)),
        "dust":   (RGBColor(0x4D, 0x6F, 0x82), RGBColor(0xD8, 0xE5, 0xED), RGBColor(0x1A, 0x2E, 0x3A)),
    }

    # ── Overview slide ──
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF7, 0xF3, 0xEC)

    # Central topic
    ct = slide.shapes.add_textbox(Inches(4.5), Inches(0.3), Inches(4.33), Inches(0.9))
    ct.fill.solid(); ct.fill.fore_color.rgb = RGBColor(0x3D, 0x2B, 0x1F)
    cf = ct.text_frame; cf.word_wrap = True
    cp = cf.paragraphs[0]
    _bionic_runs(cp, mindmap.get("topic", topic), 26, RGBColor(0xFF, 0xFF, 0xFF))

    # Summary
    sb = slide.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.33), Inches(0.8))
    sf = sb.text_frame; sf.word_wrap = True
    _bionic_runs(sf.paragraphs[0], mindmap.get("summary", ""), 16, RGBColor(0x3D, 0x2B, 0x1F))

    # Branch overview grid
    branches = mindmap.get("branches", [])
    cols = 4
    for bi, branch in enumerate(branches[:8]):
        col = bi % cols; row = bi // cols
        x = Inches(0.4 + col * 3.2); y = Inches(2.4 + row * 2.2)
        bname = branch.get("color", "teal")
        accent, bg, text = BRANCH_COLORS.get(bname, BRANCH_COLORS["teal"])
        box = slide.shapes.add_textbox(x, y, Inches(3.0), Inches(1.8))
        box.fill.solid(); box.fill.fore_color.rgb = bg
        bf = box.text_frame; bf.word_wrap = True
        bp = bf.paragraphs[0]
        bp.text = f"{branch.get('emoji', '◉')}  {branch.get('label', '')}"
        bp.font.bold = True; bp.font.size = Pt(16); bp.font.color.rgb = accent
        for sn in branch.get("subnodes", [])[:2]:
            sp = bf.add_paragraph()
            sp.text = f"  • {sn}"; sp.font.size = Pt(11); sp.font.color.rgb = text

    # Fun facts strip
    facts = mindmap.get("fun_facts", [])
    if facts:
        ffb = slide.shapes.add_textbox(Inches(0), Inches(6.8), Inches(13.33), Inches(0.55))
        ffb.fill.solid(); ffb.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xC0)
        fff = ffb.text_frame
        ffp = fff.paragraphs[0]
        ffp.text = "⚡  " + "  ·  ".join(facts[:3])
        ffp.font.size = Pt(11); ffp.font.color.rgb = RGBColor(0x4E, 0x34, 0x00)

    # ── One slide per branch ──
    for branch in branches:
        bname = branch.get("color", "teal")
        accent, bg, text = BRANCH_COLORS.get(bname, BRANCH_COLORS["teal"])

        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = bg

        # Top bar
        top = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.3))
        top.fill.solid(); top.fill.fore_color.rgb = accent; top.line.fill.background()

        # Branch heading
        hb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12.0), Inches(1.0))
        hf = hb.text_frame; hf.word_wrap = True
        hp = hf.paragraphs[0]
        hp.text = f"{branch.get('emoji', '')}  {branch.get('label', '')}"
        hp.font.bold = True; hp.font.size = Pt(28); hp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Subnodes as bionic cards
        subnodes = branch.get("subnodes", [])
        for si, sn in enumerate(subnodes[:4]):
            col = si % 2; row = si // 2
            x = Inches(0.6 + col * 6.3); y = Inches(1.6 + row * 2.4)
            snb = slide.shapes.add_textbox(x, y, Inches(5.9), Inches(2.0))
            snb.fill.solid(); snb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            snf = snb.text_frame; snf.word_wrap = True
            _bionic_runs(snf.paragraphs[0], sn, 20, text)

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


# ── Student pipeline endpoints ─────────────────────────────────────────────

@router.post("/student-adhd-slides")
async def generate_student_adhd_slides(req: AgentGenerateRequest):
    """Slide Sorter (Student ADHD): bionic reading, fun facts, mind map hints."""
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    job_id = str(uuid.uuid4())
    slides = await run_in_threadpool(agent_service.generate_student_adhd_slides, pockets, req.lang)
    topic  = pockets.get("topic", "Slides")
    pptx   = await run_in_threadpool(_build_student_adhd_pptx, slides, topic)

    artifact_cache.store(
        job_id, pptx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        f"adhd_student_{job_id[:8]}.pptx",
    )
    return {"job_id": job_id, "slides": slides, "download_url": f"/api/download/{job_id}",
            "topic": topic, "slide_count": len(slides)}


@router.post("/student-autism-slides")
async def generate_student_autism_slides(req: AgentGenerateRequest):
    """Slide Sorter (Student Autism): bionic reading, clean predictable layout."""
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    job_id = str(uuid.uuid4())
    slides = await run_in_threadpool(agent_service.generate_student_autism_slides, pockets, req.lang)
    topic  = pockets.get("topic", "Slides")
    pptx   = await run_in_threadpool(_build_student_autism_pptx, slides, topic)

    artifact_cache.store(
        job_id, pptx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        f"autism_student_{job_id[:8]}.pptx",
    )
    return {"job_id": job_id, "slides": slides, "download_url": f"/api/download/{job_id}",
            "topic": topic, "slide_count": len(slides)}


@router.post("/visual-mindmap")
async def generate_visual_mindmap(req: AgentGenerateRequest):
    """Visualizer Agent: mind map structure + fun facts + PPTX download."""
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    job_id  = str(uuid.uuid4())
    mindmap = await run_in_threadpool(agent_service.generate_visual_mindmap, pockets, req.lang)
    topic   = pockets.get("topic", "Mind Map")
    pptx    = await run_in_threadpool(_build_mindmap_pptx, mindmap, topic)

    artifact_cache.store(
        job_id, pptx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        f"mindmap_{job_id[:8]}.pptx",
    )
    return {"job_id": job_id, "mindmap": mindmap, "download_url": f"/api/download/{job_id}",
            "topic": topic}


@router.post("/examiner-worksheet")
async def generate_examiner_worksheet(req: AgentGenerateRequest):
    """Examiner Agent: self-study worksheet with fill-blanks, T/F, matching, challenge."""
    pockets = pocket_cache.get(req.file_id)
    if not pockets:
        text = text_cache.get(req.file_id)
        if not text:
            raise HTTPException(404, f"File '{req.file_id}' not found.")
        pockets = await run_in_threadpool(agent_service.read_document, text)
        pocket_cache.store(req.file_id, pockets)

    try:
        worksheet = await run_in_threadpool(agent_service.generate_examiner_worksheet, pockets, req.lang)
    except Exception as exc:
        raise HTTPException(422, "Worksheet generation failed — please retry") from exc
    return {"worksheet": worksheet, "topic": pockets.get("topic", "")}
