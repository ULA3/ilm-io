"""
Generation router — all content-creation endpoints.

POST /generate/slides
POST /generate/audiobook
POST /generate/visual-summary
POST /generate/worksheet
POST /generate/simplified-text
POST /generate/educator/{kind}

All endpoints accept {file_id, condition?} and return typed responses.
Heavy AI work runs via run_in_threadpool to avoid blocking the event loop.
Generated artifacts (PPTX, PDF, MP3) are uploaded to Supabase Storage and
the download URL is returned.
"""
from __future__ import annotations
import io
import uuid
from datetime import datetime

from fastapi import APIRouter, HTTPException
from fastapi.concurrency import run_in_threadpool
from pydantic import BaseModel

from models.schemas import (
    GenerateSlidesResponse, AudiobookResponse, VisualSummaryResponse,
    WorksheetResponse, SimplifiedTextResponse, EducatorGenerateResponse,
    Slide, VisualCard, WorksheetQuestion,
    OutputKind, EducatorOutputKind, Condition,
)
from services import ai_service, tts_service, supabase_client, text_cache, artifact_cache
from services.text_limits import TEXT_LIMIT, prepare_document_text

router = APIRouter(prefix="/generate", tags=["generate"])


# ── Request models ────────────────────────────────────────────────────────

class GenerateRequest(BaseModel):
    file_id: str
    condition: str = "general"
    difficulty: int = 2          # 1-3, used by worksheet


class EducatorGenerateRequest(BaseModel):
    file_id: str
    condition: str = "general"


class FinancialLiteracyRequest(BaseModel):
    file_id: str
    condition: str = "general"
    lang: str = "en"


class ExplainStyleRequest(BaseModel):
    file_id: str
    style: str = "simple"      # simple | story | steps | bullets | malaysian
    condition: str = "general"
    lang: str = "en"


# ── Helpers ───────────────────────────────────────────────────────────────

async def _get_text(file_id: str) -> str:
    # Fast path: in-memory cache (populated on upload)
    cached = text_cache.get(file_id)
    if cached:
        return cached
    # Slow path: Supabase (needed if server restarted between upload and generate)
    record = await run_in_threadpool(supabase_client.get_upload_sync, file_id)
    if not record:
        raise HTTPException(404, f"File '{file_id}' not found. Please upload your file again.")
    return record["text_content"]


def _save_artifact(job_id: str, kind: str, content_bytes: bytes, ext: str, content_type: str) -> str:
    """Store generated file in memory cache and return its download path."""
    filename = f"{kind}_{job_id[:8]}.{ext}"
    artifact_cache.store(job_id, content_bytes, content_type, filename)
    return f"/api/download/{job_id}"


# ── PPTX builder ──────────────────────────────────────────────────────────

def _build_pptx(slides: list[dict]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # earthy background
    BG = RGBColor(0xF7, 0xF3, 0xEC)    # cream
    TITLE_C = RGBColor(0x3D, 0x2B, 0x1F)
    BODY_C = RGBColor(0x7A, 0x5C, 0x4A)
    ACCENT = RGBColor(0x8F, 0xA6, 0x8E)

    for s in slides:
        slide = prs.slides.add_slide(blank_layout)

        # background fill
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG

        # Title
        txb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1.2))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = s.get("title", "")
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = TITLE_C

        # Bullets
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(10), Inches(4.5))
        tf2 = content_box.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(s.get("bullets", [])):
            para = tf2.add_paragraph() if i else tf2.paragraphs[0]
            para.text = f"• {bullet}"
            para.font.size = Pt(22)
            para.font.color.rgb = BODY_C
            para.space_before = Pt(8)

        # Visual hint (bottom right, small)
        hint = s.get("visual_hint", "")
        if hint:
            hb = slide.shapes.add_textbox(Inches(9), Inches(6.4), Inches(4), Inches(0.8))
            hf = hb.text_frame
            hf.paragraphs[0].text = f"Visual: {hint}"
            hf.paragraphs[0].font.size = Pt(11)
            hf.paragraphs[0].font.color.rgb = ACCENT

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── PDF builder (high-contrast, dyslexia-friendly) ────────────────────────

def _build_pdf(sections: list[dict]) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=2.5*cm, rightMargin=2.5*cm,
        topMargin=2.5*cm, bottomMargin=2.5*cm,
    )
    styles = getSampleStyleSheet()
    heading_style = ParagraphStyle(
        "IlmHeading", parent=styles["Heading1"],
        fontSize=18, leading=24,
        textColor=HexColor("#3D2B1F"),
        spaceAfter=8,
    )
    body_style = ParagraphStyle(
        "IlmBody", parent=styles["Normal"],
        fontSize=14, leading=22,
        textColor=HexColor("#7A5C4A"),
        fontName="Helvetica",
        spaceAfter=6,
    )
    key_style = ParagraphStyle(
        "IlmKey", parent=styles["Normal"],
        fontSize=12, leading=18,
        textColor=HexColor("#5E7D5C"),
        fontName="Helvetica-Bold",
        spaceAfter=16,
        leftIndent=12,
    )

    story = []
    for sec in sections:
        story.append(Paragraph(sec.get("heading", ""), heading_style))
        story.append(Paragraph(sec.get("text", ""), body_style))
        if sec.get("key_idea"):
            story.append(Paragraph(f"Key idea: {sec['key_idea']}", key_style))
        story.append(Spacer(1, 0.3*cm))

    doc.build(story)
    return buf.getvalue()


# ── Endpoints ─────────────────────────────────────────────────────────────

@router.post("/slides", response_model=GenerateSlidesResponse)
async def gen_slides(req: GenerateRequest):
    text = await _get_text(req.file_id)
    truncated = len(text) > TEXT_LIMIT
    job_id = str(uuid.uuid4())
    raw_slides = await run_in_threadpool(ai_service.generate_slides, text, req.condition)
    pptx_bytes = await run_in_threadpool(_build_pptx, raw_slides)
    url = _save_artifact(job_id, "slides", pptx_bytes, "pptx",
                         "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    slides = [Slide(**s) for s in raw_slides]
    return GenerateSlidesResponse(job_id=job_id, kind=OutputKind.slides,
                                  slides=slides, download_url=url, condition=req.condition,
                                  truncated=truncated)


@router.post("/audiobook", response_model=AudiobookResponse)
async def gen_audiobook(req: GenerateRequest):
    text = await _get_text(req.file_id)
    job_id = str(uuid.uuid4())
    # Limit text to ~3000 chars to stay within free-tier limits
    tts_text = text[:3000]
    mp3_bytes = await tts_service.text_to_speech(tts_text)
    url = _save_artifact(job_id, "audio", mp3_bytes, "mp3", "audio/mpeg")
    duration = max(1, len(mp3_bytes) // 16000)
    return AudiobookResponse(job_id=job_id, audio_url=url,
                             duration_seconds=duration,
                             transcript=tts_text[:400] + "…" if len(tts_text) > 400 else tts_text)


@router.post("/visual-summary", response_model=VisualSummaryResponse)
async def gen_visual_summary(req: GenerateRequest):
    text = await _get_text(req.file_id)
    truncated = len(text) > TEXT_LIMIT
    job_id = str(uuid.uuid4())
    raw_cards = await run_in_threadpool(ai_service.generate_visual_summary, text)
    cards = [VisualCard(**c) for c in raw_cards]
    sections = [{"heading": c.concept, "text": f"{c.explanation}\n\nAnalogy: {c.analogy}", "key_idea": ""} for c in cards]
    pdf_bytes = await run_in_threadpool(_build_pdf, sections)
    url = _save_artifact(job_id, "visual", pdf_bytes, "pdf", "application/pdf")
    return VisualSummaryResponse(job_id=job_id, cards=cards, download_url=url, truncated=truncated)


@router.post("/worksheet", response_model=WorksheetResponse)
async def gen_worksheet(req: GenerateRequest):
    text = await _get_text(req.file_id)
    truncated = len(text) > TEXT_LIMIT
    job_id = str(uuid.uuid4())
    ws = await run_in_threadpool(ai_service.generate_worksheet, text, req.condition, req.difficulty)
    questions = [WorksheetQuestion(**q) for q in ws["questions"]]
    sections = [{"heading": f"Q{q.number}. {q.question}",
                 "text": "\n".join(q.options) if q.options else "",
                 "key_idea": q.answer} for q in questions]
    pdf_bytes = await run_in_threadpool(_build_pdf, sections)
    url = _save_artifact(job_id, "worksheet", pdf_bytes, "pdf", "application/pdf")
    return WorksheetResponse(job_id=job_id, title=ws.get("title", "Worksheet"),
                             questions=questions, download_url=url,
                             adapted_for=ws.get("adapted_for", req.condition),
                             truncated=truncated)


@router.post("/simplified-text", response_model=SimplifiedTextResponse)
async def gen_simplified(req: GenerateRequest):
    text = await _get_text(req.file_id)
    truncated = len(text) > TEXT_LIMIT
    job_id = str(uuid.uuid4())
    result = await run_in_threadpool(ai_service.generate_simplified_text, text)
    sections = result["sections"]
    pdf_bytes = await run_in_threadpool(_build_pdf, sections)
    url = _save_artifact(job_id, "simplified", pdf_bytes, "pdf", "application/pdf")
    return SimplifiedTextResponse(job_id=job_id, original_word_count=len(text.split()),
                                  simplified_word_count=sum(len(s["text"].split()) for s in sections),
                                  sections=sections, download_url=url, truncated=truncated)


@router.post("/educator/{kind}", response_model=EducatorGenerateResponse)
async def gen_educator(kind: str, req: EducatorGenerateRequest):
    valid_kinds = {e.value for e in EducatorOutputKind}
    if kind not in valid_kinds:
        raise HTTPException(400, f"Unknown kind '{kind}'.")
    text = await _get_text(req.file_id)
    truncated = len(text) > TEXT_LIMIT
    job_id = str(uuid.uuid4())
    content = await run_in_threadpool(ai_service.generate_educator_content, text, kind, req.condition)
    sections = [{"heading": k.replace("_", " ").title(), "text": str(v), "key_idea": ""} for k, v in content.items()]
    pdf_bytes = await run_in_threadpool(_build_pdf, sections)
    url = _save_artifact(job_id, kind, pdf_bytes, "pdf", "application/pdf")
    return EducatorGenerateResponse(job_id=job_id, kind=EducatorOutputKind(kind),
                                    condition=Condition(req.condition),
                                    content=content, download_url=url, truncated=truncated)


@router.post("/explain-style")
async def gen_explain_style(req: ExplainStyleRequest):
    valid_styles = {"simple", "story", "steps", "bullets", "malaysian"}
    if req.style not in valid_styles:
        raise HTTPException(400, f"Unknown style '{req.style}'. Valid: {sorted(valid_styles)}")
    text = await _get_text(req.file_id)
    truncated = len(text) > TEXT_LIMIT
    result = await run_in_threadpool(
        ai_service.generate_with_style, text, req.style, req.condition, req.lang
    )
    return {
        "style":     result.get("style", req.style),
        "title":     result.get("title", ""),
        "content":   result.get("content", ""),
        "condition": req.condition,
        "truncated": truncated,
    }


@router.post("/financial-literacy")
async def gen_financial_literacy(req: FinancialLiteracyRequest):
    text = await _get_text(req.file_id)
    truncated = len(text) > TEXT_LIMIT
    job_id = str(uuid.uuid4())
    result = await run_in_threadpool(
        ai_service.generate_financial_literacy, text, req.condition, req.lang
    )
    return {
        "job_id":     job_id,
        "title":      result.get("title", "Financial Literacy Lesson"),
        "concepts":   result.get("concepts", []),
        "condition":  req.condition,
        "truncated":  truncated,
    }
